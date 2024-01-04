import openai
import json

# optional; defaults to `os.environ['OPENAI_API_KEY']`
openai.api_key = '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'

messages=[
    {
        "role": "user",
        "content": "Generate a PPT related to Marie Curie",
    },
]
slide_number="1"
prompt="Generate a PPT of 3 slides related to Albert Einstein"
word_content=""
slideDeck={"slide_number": float(slide_number), "title": "", "content": "", "narration": ""}

messages = [
        {"role": "system", "content": "User will ask you to create or update text content for some slides"+(" based on the aforementioned Input Article" if word_content else "")+". The response format should be a valid json format structured as this: [{\"slide_number\": <Float>, \"title\": \"<String>\", \"content\": \"<String>\", \"narration\": \"<String>\"},{\"slide_number\": <Float>, \"title\": \"<String>\", \"content\": \"<String>\", \"narration\": \"<String>\"}] \n content field in the response comprehensive enough as it is the main text of each slide. \n For content use a mix of bullet points and text when applicable. \n If you are modifying an existing slide leave the slide number unchanged but if you are adding slides to the existing slides, use decimal digits for the slide number. for example to add a slide after slide 2, use slide number 2.1, 2.2, ... \n If user asks to remove a slide, set its slide number to negative of its current value because slides with negative slide number will be excluded from presentation. \n The existing slides are as follows: "+json.dumps(slideDeck)},
        {"role": "system", "content": "For each slide the content field is the main body of the slide while the narration field is just an example transcript of the presentation of the content field. \n Never mention the slide number in the transcript."},
        {"role": "system", "content": "For each slide, the content field should be the default field to modify if modification is demanded by the user for the slide, not the narration field. "},
        {"role": "system", "content": "For each slide, the narration field should only be populated if explicitely asked in user prompt, otherwise should be left empty. "},
        {"role": "system", "content": "Response should be valid json in the format described earlier. slide_number, title,and content are mandatory keys."},
        {"role": "user", "content": prompt}
    ]
completion = openai.chat.completions.create(
    model="gpt-3.5-turbo",
    messages=messages,
)
# 
slides_data = json.loads(completion.choices[0].message.content)

print(slides_data)
print(type(slides_data))
from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# # Updated slide data
# slides_data = [
#     {"slide_number": 1.0, "title": "Introduction to Marie Curie", "content": "Marie Curie was a physicist and chemist who conducted groundbreaking research on radioactivity. She was the first woman to win a Nobel Prize, the only person to win the Nobel Prize in two different sciences, and the first woman to become a professor at the University of Paris.", "narration": ""},
#     {"slide_number": 2.0, "title": "Early Life", "content": "Born as Maria Sklodowska on November 7, 1867, in Warsaw, Poland. Her parents were both teachers, and she was the youngest of five children.", "narration": ""},
#     {"slide_number": 3.0, "title": "Education", "content": "In 1891, she moved to Paris to continue her studies at the Sorbonne where she earned degrees in physics and mathematical sciences.", "narration": ""},
#     {"slide_number": 4.0, "title": "Marriage and Research", "content": "She met Pierre Curie, professor in the School of Physics in 1894 and in the following year, they were married. They conducted research on radioactive substances.", "narration": ""},
#     {"slide_number": 5.0, "title": "Discovering New Elements", "content": "The Curies and their student, Gustave Bémont, conducted a series of experiments which led to the discovery of polonium and radium.", "narration": ""},
#     {"slide_number": 6.0, "title": "Nobel Prizes", "content": "Curie became the first woman to win the Nobel Prize in Physics in 1903. She won another in 1911, this time in Chemistry.", "narration": ""},
#     {"slide_number": 7.0, "title": "Pioneering Research in Radioactivity", "content": "Curie dedicated her life to the study of radioactivity and through her research, she established a foundation for modern atomic physics.", "narration": ""},
#     {"slide_number": 8.0, "title": "World War I", "content": "During World War I, Curie operated mobile radiography units, also known as Little Curies, to provide X-ray services to field hospitals.", "narration": ""},
#     {"slide_number": 9.0, "title": "Death and Legacy", "content": "Marie Curie died on July 4, 1934, from aplastic anemia, believed to be caused by prolonged exposure to radiation. Her legacy lives on in the numerous awards carrying her name and in the Marie Curie Actions for researchers.", "narration": ""},
#     {"slide_number": 10.0, "title": "Conclusion", "content": "Marie Curie's dedication to science and humanity has had a profound impact on the world. Her life’s work has opened up new paths for research in physics and medicine, and she has left an enduring legacy.", "narration": ""}
# ]

# Function to add a slide
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Using layout 1 for title and content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    tf = body_placeholder.text_frame
    tf.text = content

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)

# Add slides
for slide in slides_data:
    add_slide(slide['title'], slide['content'])

# Save the presentation
prs.save('Output_Presentation.pptx')
