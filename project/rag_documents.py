import streamlit as st
import os
from pathlib import Path, PurePath
from langchain.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import DirectoryLoader
from langchain.chains import RetrievalQA
from langchain.prompts import PromptTemplate
from langchain.vectorstores import FAISS
from cg_utils import *
import re
from pptx import Presentation

# Directories
UPLOAD_DIR = Path(__file__).resolve().parent.joinpath('', 'upload')
INPUT_DIR = Path(__file__).resolve().parent.joinpath('', 'input')
VECTOR_STORE_DIR = Path(__file__).resolve().parent.joinpath('', 'vector_store')
if not os.path.exists(UPLOAD_DIR):
   os.makedirs(UPLOAD_DIR)
if not os.path.exists(INPUT_DIR):
   os.makedirs(INPUT_DIR)   
if not os.path.exists(VECTOR_STORE_DIR):
   os.makedirs(VECTOR_STORE_DIR)


# Get text-to-text FMs
t2t_fms = get_t2t_fms(fm_vendors)


def process_docs(in_dir:str, out_dir:str):
    """Save uploaded files, process files for embeddings and move processed files"""
    with st.spinner('Splitting files, generating and storing embeddings...'):
        for doc in st.session_state.rag_docs_key:
            upload_path = Path(UPLOAD_DIR, doc.name)
            with open(upload_path, mode='wb') as w:
                w.write(doc.getvalue())
        pdf_chunks = split_pdfs(in_dir)
        st.session_state.faiss_vector_db = embeddings_faiss(pdf_chunks,VECTOR_STORE_DIR, bedrock_embeddings)                
        for f in os.listdir(in_dir):
            src_path = os.path.join(in_dir, f)
            dst_path = os.path.join(out_dir, f)
            os.rename(src_path, dst_path)
        
   
def list_files(dir:str):
    """Display file listing for a directory"""
    if len(os.listdir(dir)) > 0:
        for i,f in enumerate(os.listdir(dir)):
            st.markdown(f"-  {f}",unsafe_allow_html=True)


def split_pdfs(pdf_dir:str) -> list:
    """Load PDFs from a directory and split PDFs into tokens"""
    if os.listdir(pdf_dir):
        loader = DirectoryLoader(pdf_dir, loader_cls=PyPDFLoader)
        documents = loader.load()
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=0,
            separators=["\n", "\n\n", "(?<=\. )"]
            )
        document_chunks = text_splitter.split_documents(documents)
        return document_chunks


def embeddings_faiss(doc_list:list, db_dir:str, embeddings_fn):
    """Generate embeddings for document chunks and store in a FAISS vector datastore"""
    if doc_list:
        new_db = FAISS.from_documents(doc_list, embeddings_fn)
        if os.listdir(db_dir):
            local_db = FAISS.load_local(db_dir, embeddings_fn)
            local_db.merge_from(new_db)
            local_db.save_local(db_dir)
            return local_db
        else:
            new_db.save_local(db_dir)
            return new_db
    elif os.listdir(db_dir):
        local_db = FAISS.load_local(db_dir, embeddings_fn)
        return local_db
    else:
        return None


def empty_dir(data_dir:str):
    """Remove files from a dircetory"""
    for f in os.listdir(data_dir):
        os.remove(os.path.join(data_dir, f))
      

def ask_fm_rag_off(prompt:str, modelid:str):
    """FM query - RAG disabled"""
    return ask_fm(modelid, prompt)

# Function that I need!
def ask_fm_rag_on(prompt:str, modelid:str, vector_db):
    """FM contextual query - RAG enabled"""
    prompt_template = """
    Human: Use only the following context to provide a concise answer to the question at the end. 
    If you cannot find the answer in the context, just say that you don't know, don't try to make up an answer.

    {context}

    Question: {question}
    Assistant:
    """             
    PROMPT = PromptTemplate(
        template=prompt_template, input_variables=["context", "question"]
    )
    qa = RetrievalQA.from_chain_type(
        llm=get_fm(modelid),
        chain_type="stuff",
        retriever=vector_db.as_retriever(
            search_type="similarity",
            search_kwargs={"k": 5}
        ),
        return_source_documents=True,
        chain_type_kwargs={"prompt": PROMPT}
    )
    result = qa({"query": prompt})
    if result['source_documents'] == []:
        response = f"""{result['result']}"""
    else:
        response = f"""{result['result']}<br /><br /> <b>Source:</b> {PurePath(result['source_documents'][0].metadata['source']).name}"""
    print(result['result'])
    # return response
    return result['result']

# rag_fm_prompt = "WHAT IS Wheel Alignment"
# rag_fm = "anthropic.claude-v2:1"
# #"anthropic.claude-3-sonnet-20240229-v1:0"
# ask_fm_rag_on(rag_fm_prompt, rag_fm, embeddings_faiss(split_pdfs(INPUT_DIR),VECTOR_STORE_DIR, bedrock_embeddings))



def read_prompt_from_ppt_template(ppt_template_path, generated_ppt_path, rag_fm, faiss_vector_db):

    # Load the presentation
    ppt = Presentation("template/template1.pptx")
    pattern = r'\[.*?\]'
    
    # Iterate through slides and then through text frames to replace placeholders
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    # Concatenate all runs in the paragraph to get the full paragraph text
                    full_text = ''.join([run.text for run in paragraph.runs])
                    
                    # Search for patterns and do something with them (e.g., print or replace)
                    matches = re.findall(pattern, full_text)
                    for match in matches:
                        print(match)
                        full_text = full_text.replace(match, ask_fm_rag_on(match, rag_fm, faiss_vector_db))
                    paragraph.clear()  # Clear the existing runs
                    paragraph.add_run().text = full_text  # Add new run with the updated text
    # Save the modified PPTX file
    ppt.save(generated_ppt_path)