import re
from pptx import Presentation

def read_prompt_from_ppt_template(ppt_template_path, generated_ppt_path):

    # Load the presentation
    ppt = Presentation("template/template1.pptx")
    pattern = r'\[.*?\]'
    
    # Iterate through slides and then through text frames to replace placeholders
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    # Concatenate all runs in the paragraph to get the full paragraph text
                    full_text = ''.join([run.text for run in paragraph.runs])
                    
                    # Search for patterns and do something with them (e.g., print or replace)
                    matches = re.findall(pattern, full_text)
                    for match in matches:
                        print(match)
                        full_text = full_text.replace(match, ask_fm_rag_on(match, rag_fm, embeddings_faiss(split_pdfs(INPUT_DIR),VECTOR_STORE_DIR, bedrock_embeddings)))
                    paragraph.clear()  # Clear the existing runs
                    paragraph.add_run().text = full_text  # Add new run with the updated text
    # Save the modified PPTX file
    ppt.save(generated_ppt_path)