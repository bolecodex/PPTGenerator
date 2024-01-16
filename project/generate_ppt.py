from pptx import Presentation
from pptx.util import Inches, Pt

# Updated slide data
slides_data = [
    {
        "slide_number": 1.0,
        "title": "Introduction to Albert Einstein",
        "content": "Albert Einstein was born on March 14, 1879 in Ulm, in the Kingdom of Württemberg in the German Empire. He is best known for developing the theory of relativity, which revolutionized the understanding of space, time, and gravity.",
        "narration": ""
    },
    {
        "slide_number": 2.0,
        "title": "Einstein's Contributions to Science",
        "content": "• E=mc^2: The mass-energy equivalence.\n• Theory of General relativity: This led to prediction of the deflection of light by gravity, concept of black holes and Big Bang theory.\n• Quantum Theory: Einstein made important contributions to early quantum theory.",
        "narration": ""
    },
    {
        "slide_number": 3.0,
        "title": "Einstein's Legacy",
        "content": "Einstein's work continues to influence the course of science to this day. His theories have been instrumental in developing GPS technology and understanding the universe. He was awarded the Nobel Prize in Physics in 1921.",
        "narration": ""
    }
]

def generate_ppt(slides_data, output_file_name):
    """
    Creates a PowerPoint presentation from the given slides data.

    Args:
    slides_data (list of dict): A list of dictionaries, each representing a slide's content.
    output_file_name (str): The name of the file to save the presentation.

    Returns:
    Presentation: The PowerPoint presentation object.
    """
    prs = Presentation()

    def add_slide(prs, title, content):
        """
        Adds a slide to the presentation with the given title and content.

        Args:
        prs (Presentation): The PowerPoint presentation object.
        title (str): The title of the slide.
        content (str): The content of the slide.
        """
        slide_layout = prs.slide_layouts[1]  # Using layout 1 for title and content
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        body_placeholder = slide.placeholders[1]

        title_placeholder.text = title
        tf = body_placeholder.text_frame
        tf.text = content

        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)

    # Add slides based on the provided data
    for slide in slides_data:
        add_slide(prs, slide['title'], slide['content'])

    # Save the presentation
    prs.save(output_file_name)
    return prs

generate_ppt(slides_data, 'Albert_Einstein_Presentation.pptx')
