from pptx import Presentation
import json

def ppt_to_json(ppt_file):
    prs = Presentation(ppt_file)
    slides_data = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        # Initialize the slide data dictionary
        slide_data = {"slide_number": float(slide_number), "title": "", "content": "", "narration": ""}

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            if shape == slide.shapes.title:
                slide_data["title"] = shape.text
            else:
                slide_data["content"] += shape.text + "\n"

        slides_data.append(slide_data)

    return json.dumps(slides_data)

# Example usage
# ppt_file = 'Marie_Curie_Presentation.pptx'
# slides_json = ppt_to_json(ppt_file)
# # print(slides_json)
# print(json.dumps(slides_json, indent=4))


def get_latest_slide_deck(pptx_file_path, current_slide_deck, logger):
    """
    Retrieves information from the latest PowerPoint slide deck in the specified file.

    Args:
        pptx_file_path (str): The path to the PowerPoint file.
        current_slide_deck (list): The current list of slide information.
        logger (Logger): The logger object for error reporting.

    Returns:
        list: A list of dictionaries containing slide information, including slide number,
              title, content, and narration.

    Note:
        This function attempts to load the PowerPoint presentation and extract slide
        information. If the file is open in another application with a read lock (rare), it will prompt the user
        to close it and retry. If the file does not exist, it returns the current slide deck.

    """
    presentation = None
    while True:
        try:
            presentation = Presentation(pptx_file_path)
            break
        except PermissionError:
            # print(format_warning("Action Required:")+ f" The file {output_file} is open in PowerPoint or another application." +format_prompt("Please close it and press Enter to retry. "))
            # logger.error(f"Error: The file '{output_file}' is open in PowerPoint or another application. Please close it and press Enter to retry.")
            input()  # Wait for the user to close the file
        except FileNotFoundError:
            return current_slide_deck           
        except Exception as e:
            # logger.error(f"An unexpected error occurred while loading the presentation: {e}")
            return current_slide_deck            

    updated_slide_deck = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_info = {"slide_number": float(slide_number), "title": "", "content": "", "narration": ""}
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape == slide.shapes[0]:  # Assuming the first shape is the title
                    slide_info["title"] = shape.text
                else:
                    slide_info["content"] += shape.text + "\n"
        # Extract notes from the notes slide
        notes_slide = slide.notes_slide
        for shape in notes_slide.shapes:
            if shape.has_text_frame:
                slide_info["narration"] += shape.text + "\n"

        updated_slide_deck.append(slide_info)

    return updated_slide_deck

